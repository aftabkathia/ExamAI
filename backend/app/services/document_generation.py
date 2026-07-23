"""Generate PDF, Word, Excel, and PowerPoint files from chat requests."""

from __future__ import annotations

import base64
import io
import json
import re
from typing import Any, Literal

from groq import Groq

TEXT_MODEL = "llama-3.3-70b-versatile"
DocType = Literal["pdf", "docx", "xlsx", "pptx"]

_CREATE = re.compile(
    r"\b(create|make|generate|export|build|prepare|write|produce|banao|banado|bana\s*do|tyar)\b",
    re.IGNORECASE,
)

_TYPE_HINTS: list[tuple[DocType, re.Pattern[str]]] = [
    ("pdf", re.compile(r"\b(pdf|portable\s+document)\b", re.I)),
    ("docx", re.compile(r"\b(word(\s+document)?|docx?)\b", re.I)),
    ("xlsx", re.compile(r"\b(excel(\s+sheet)?|xlsx?|spreadsheet|worksheet)\b", re.I)),
    ("pptx", re.compile(r"\b(powerpoint|pptx?|presentation|slides?)\b", re.I)),
]

_MIME: dict[DocType, str] = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_EXT: dict[DocType, str] = {
    "pdf": "pdf",
    "docx": "docx",
    "xlsx": "xlsx",
    "pptx": "pptx",
}

_JSON_PROMPTS: dict[DocType, str] = {
    "pdf": """Create document content as JSON for a PDF. Schema:
{"title":"string","sections":[{"heading":"string","paragraphs":["string"]}],"bullets":["optional list items"]}
Use clear exam-quality content. 3-8 sections. Output ONLY valid JSON.""",
    "docx": """Create document content as JSON for a Word file. Schema:
{"title":"string","sections":[{"heading":"string","paragraphs":["string"]}],"bullets":["optional"]}
Use clear exam-quality content. 3-8 sections. Output ONLY valid JSON.""",
    "xlsx": """Create spreadsheet content as JSON. Schema:
{"title":"string","sheets":[{"name":"Sheet1","headers":["Col1","Col2"],"rows":[["r1c1","r1c2"]]}]}
Include useful tabular data (MCQs, dates, formulas, stats). 1-3 sheets, up to 30 rows each. Output ONLY valid JSON.""",
    "pptx": """Create presentation content as JSON. Schema:
{"title":"string","slides":[{"title":"string","bullets":["point1","point2"]}]}
5-12 slides for teaching/revision. Output ONLY valid JSON.""",
}


def detect_document_type(message: str, has_uploads: bool = False) -> DocType | None:
    if has_uploads:
        return None
    text = (message or "").strip()
    if len(text) < 6:
        return None

    for doc_type, pattern in _TYPE_HINTS:
        if not pattern.search(text):
            continue
        if _CREATE.search(text) or re.search(
            r"\b(export|download|save|as|in)\s+(a\s+)?(pdf|word|excel|ppt|docx?|xlsx?|pptx?|presentation|spreadsheet)\b",
            text,
            re.I,
        ):
            return doc_type

    # "document/file banao" without type → default PDF
    if _CREATE.search(text) and re.search(
        r"\b(document|file|notes\s+file|study\s+material)\b", text, re.I
    ):
        return "pdf"
    return None


def wants_document_generation(message: str, has_uploads: bool = False) -> bool:
    return detect_document_type(message, has_uploads) is not None


def _parse_json(raw: str) -> dict[str, Any]:
    text = (raw or "").strip()
    text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.IGNORECASE)
    start = text.find("{")
    end = text.rfind("}")
    if start >= 0 and end > start:
        text = text[start : end + 1]
    return json.loads(text)


def _content_from_groq(client: Groq, user_message: str, doc_type: DocType) -> dict[str, Any]:
    response = client.chat.completions.create(
        model=TEXT_MODEL,
        messages=[
            {"role": "system", "content": _JSON_PROMPTS[doc_type]},
            {"role": "user", "content": user_message},
        ],
        temperature=0.35,
        max_tokens=2500,
        response_format={"type": "json_object"},
    )
    return _parse_json(response.choices[0].message.content or "{}")


def _safe_name(title: str, ext: str) -> str:
    base = re.sub(r"[^\w\s-]", "", title or "document").strip().replace(" ", "_")
    base = base[:48] or "examai_document"
    return f"{base}.{ext}"


def _pdf_text(value: Any) -> str:
    text = str(value).replace("•", "- ").replace("—", "-").replace("–", "-")
    return text.encode("latin-1", errors="replace").decode("latin-1")


def _build_pdf(data: dict[str, Any]) -> bytes:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    w = pdf.epw
    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(w, 10, _pdf_text(data.get("title", "ExamAI Document")))
    pdf.ln(4)
    for sec in data.get("sections", []):
        pdf.set_font("Helvetica", "B", 13)
        pdf.multi_cell(w, 8, _pdf_text(sec.get("heading", "Section")))
        pdf.ln(1)
        pdf.set_font("Helvetica", "", 11)
        for para in sec.get("paragraphs", []):
            line = _pdf_text(para).replace("\n", " ").strip()
            if line:
                pdf.multi_cell(w, 6, line)
                pdf.ln(1)
    bullets = data.get("bullets") or []
    if bullets:
        pdf.set_font("Helvetica", "B", 13)
        pdf.multi_cell(w, 8, "Key Points")
        pdf.ln(1)
        pdf.set_font("Helvetica", "", 11)
        for item in bullets:
            line = _pdf_text(item).replace("\n", " ").strip()
            if line:
                pdf.multi_cell(w, 6, f"- {line}")
                pdf.ln(1)
    return bytes(pdf.output())


def _build_docx(data: dict[str, Any]) -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading(data.get("title", "ExamAI Document"), 0)
    for sec in data.get("sections", []):
        doc.add_heading(sec.get("heading", "Section"), level=1)
        for para in sec.get("paragraphs", []):
            doc.add_paragraph(str(para))
    if data.get("bullets"):
        doc.add_heading("Key Points", level=1)
        for item in data["bullets"]:
            doc.add_paragraph(str(item), style="List Bullet")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_xlsx(data: dict[str, Any]) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    wb.remove(wb.active)
    sheets = data.get("sheets") or [{"name": "Sheet1", "headers": ["Topic"], "rows": [["Data"]]}]
    for i, sheet in enumerate(sheets[:3]):
        ws = wb.create_sheet(title=str(sheet.get("name", f"Sheet{i + 1}"))[:31])
        headers = sheet.get("headers") or []
        if headers:
            ws.append([str(h) for h in headers])
            for cell in ws[1]:
                cell.font = Font(bold=True)
        for row in sheet.get("rows", [])[:50]:
            ws.append([str(c) for c in row])
    if not wb.sheetnames:
        ws = wb.create_sheet("Sheet1")
        ws.append(["Title", data.get("title", "ExamAI")])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_pptx(data: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = data.get("title", "ExamAI Presentation")
    title_slide.placeholders[1].text = "Generated by ExamAI Assistant"

    for slide_data in data.get("slides", [])[:15]:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", "Slide")
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for j, bullet in enumerate(slide_data.get("bullets", [])[:8]):
            p = body.paragraphs[0] if j == 0 else body.add_paragraph()
            p.text = str(bullet)
            p.level = 0
            p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_BUILDERS = {
    "pdf": _build_pdf,
    "docx": _build_docx,
    "xlsx": _build_xlsx,
    "pptx": _build_pptx,
}


def generate_document_from_message(
    client: Groq,
    user_message: str,
    db: Any | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    doc_type = detect_document_type(user_message)
    if not doc_type:
        raise ValueError("Not a document generation request")

    content = _content_from_groq(client, user_message, doc_type)
    raw = _BUILDERS[doc_type](content)
    b64 = base64.b64encode(raw).decode("ascii")
    filename = _safe_name(content.get("title", "examai_document"), _EXT[doc_type])
    mime = _MIME[doc_type]

    label = {"pdf": "PDF", "docx": "Word", "xlsx": "Excel", "pptx": "PowerPoint"}[doc_type]
    short = user_message.strip()
    if len(short) > 80:
        short = short[:77] + "…"

    if db is not None and user_id:
        _save_to_workspace(
            db, user_id, content.get("title", filename), filename, mime, b64
        )

    return {
        "reply": (
            f"**{label} file ready** — _{content.get('title', filename)}_\n\n"
            f"Request: _{short}_\n\n"
            "Saved to your **Documents** tab. Use the download button below too."
        ),
        "generated_files": [
            {"name": filename, "mime": mime, "data_base64": b64},
        ],
        "grounded": False,
        "model": f"doc-gen-{doc_type}",
    }


def _save_to_workspace(
    db: Any, user_id: int, title: str, filename: str, mime: str, b64: str
) -> None:
    from app.services.workspace import save_document

    save_document(db, user_id, title, filename, mime, b64, source="generated")
